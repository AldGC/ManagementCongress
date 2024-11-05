from django.shortcuts import render
from django.http import HttpResponse
import os
import tempfile
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def convert_pptx_to_pdf(pptx_path, pdf_path):
    """
    Convierte un archivo PPTX a PDF usando comtypes en Windows
    """
    import comtypes.client

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    deck = powerpoint.Presentations.Open(pptx_path)
    deck.SaveAs(pdf_path, 32)  # 32 es el código para formato PDF
    deck.Close()
    powerpoint.Quit()


def convert_pptx_to_pdf_linux(pptx_path, pdf_path):
    """
    Convierte un archivo PPTX a PDF usando unoconv en Linux
    """
    import subprocess
    command = ['unoconv', '-f', 'pdf', '-o', pdf_path, pptx_path]
    subprocess.run(command)


def home(request):
    if request.method == "POST":
        nombre_completo = request.POST.get("nombre_completo")

        # Cargar y editar la presentación
        prs = Presentation("data_collection/templates/src/files/Constancia.pptx")
        slide = prs.slides[0]  # Acceder a la primera diapositiva

        nombre_encontrado = False

        # Buscar y modificar el texto del nombre
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            # Buscar el cuadro que contiene el texto "a" solo
            if text == "a":
                continue

            # El siguiente shape después del "a" será el nombre
            if not nombre_encontrado and shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.clear()

                # Crear un nuevo párrafo y establecer sus propiedades
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER

                # Agregar el texto como un nuevo run
                run = p.add_run()
                run.text = nombre_completo

                # Aplicar el formato de fuente de manera explícita
                font = run.font
                font.name = 'HelveticaNowText'
                font.size = Pt(28)
                # Color gris oscuro (RGB: 68, 68, 68)
                font.color.rgb = RGBColor(68, 68, 68)

                # Ajustar el espaciado
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                p.line_spacing = 1.0  # Espaciado de línea normal

                # Desactivar el ajuste automático de texto para mantener el tamaño de fuente
                text_frame.word_wrap = True
                text_frame.auto_size = None  # Deshabilitar auto-size

                # Ajustar el espaciado entre caracteres (tracking)
                run.font.spacing = Pt(0.5)  # Ajusta este valor según sea necesario

                nombre_encontrado = True
                break

        # Crear archivos temporales para PPTX y PDF
        temp_dir = tempfile.mkdtemp()
        temp_pptx_path = os.path.join(temp_dir, "temp.pptx")
        temp_pdf_path = os.path.join(temp_dir, "temp.pdf")

        # Guardar la presentación modificada
        prs.save(temp_pptx_path)

        try:
            # Intentar conversión en Windows
            convert_pptx_to_pdf(temp_pptx_path, temp_pdf_path)
        except:
            try:
                # Si falla, intentar conversión en Linux
                convert_pptx_to_pdf_linux(temp_pptx_path, temp_pdf_path)
            except Exception as e:
                return HttpResponse(f"Error en la conversión: {str(e)}")

        # Leer el archivo PDF generado
        with open(temp_pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="Constancia_{nombre_completo}.pdf"'

        # Limpiar archivos temporales
        try:
            os.remove(temp_pptx_path)
            os.remove(temp_pdf_path)
            os.rmdir(temp_dir)
        except:
            pass
        

        return response

    return render(request, "data_collection/home.html")